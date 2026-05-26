"""
Genera el PPT "Plan IA - Operaciones y Finanzas" para Gerencias UnionX.
27 slides.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Paleta
AZUL = RGBColor(0x4A, 0x90, 0xE2)
AZUL_OSCURO = RGBColor(0x1E, 0x40, 0xAF)
GRIS_OSCURO = RGBColor(0x0F, 0x17, 0x2A)
GRIS_MEDIO = RGBColor(0x47, 0x55, 0x69)
GRIS_CLARO = RGBColor(0xCB, 0xD5, 0xE1)
VERDE = RGBColor(0x10, 0xB9, 0x81)
VERDE_OSCURO = RGBColor(0x05, 0x96, 0x69)
ROJO = RGBColor(0xDC, 0x26, 0x26)
NARANJA = RGBColor(0xEA, 0x58, 0x0C)
AMARILLO = RGBColor(0xFA, 0xCC, 0x15)
FONDO = RGBColor(0xF8, 0xFA, 0xFC)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)

BASE = r"G:\Mi unidad\TRABAJO\RESPALDO\OPERACIONES\UNION X - IA"
LOGO = os.path.join(BASE, "data", "branding", "unionx_logo.png")
OUTPUT = os.path.join(BASE, "data", "outputs", "Plan_IA_Operaciones_Finanzas_2026-05-26.pptx")
TOTAL_SLIDES = 27

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide(bg=BLANCO):
    s = prs.slides.add_slide(BLANK)
    bg_shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg
    bg_shape.line.fill.background()
    bg_shape.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=GRIS_OSCURO,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(text, list):
        for i, line in enumerate(text):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            run = p.add_run(); run.text = line
            run.font.size = Pt(size); run.font.bold = bold
            run.font.color.rgb = color; run.font.name = font
    else:
        p = tf.paragraphs[0]; p.alignment = align
        run = p.add_run(); run.text = text
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = color; run.font.name = font
    return box


def add_rect(slide, x, y, w, h, fill=AZUL, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line; shp.line.width = Pt(0.5)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_rounded(slide, x, y, w, h, fill=BLANCO, line=GRIS_CLARO):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    shp.adjustments[0] = 0.08
    return shp


def add_footer(slide, slide_num):
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(7.15), Inches(12.833), Inches(7.15))
    line.line.color.rgb = GRIS_CLARO
    line.line.width = Pt(0.75)
    add_text(slide, 0.5, 7.2, 6, 0.3, "Plan IA · Operaciones y Finanzas · Mayo 2026", size=9, color=GRIS_MEDIO)
    slide.shapes.add_picture(LOGO, Inches(6.0), Inches(7.18), height=Inches(0.25))
    add_text(slide, 11.8, 7.2, 1.0, 0.3, f"{slide_num} / {TOTAL_SLIDES}", size=9, color=GRIS_MEDIO, align=PP_ALIGN.RIGHT)


def add_title(slide, title, subtitle=None):
    add_text(slide, 0.6, 0.5, 12.2, 0.7, title, size=28, bold=True, color=AZUL_OSCURO)
    if subtitle:
        add_text(slide, 0.6, 1.15, 12.2, 0.4, subtitle, size=14, color=GRIS_MEDIO)


# ========== SLIDE 1 — PORTADA ==========
s = add_slide(bg=GRIS_OSCURO)
add_rect(s, 0, 0, 4.5, 7.5, fill=AZUL)
s.shapes.add_picture(LOGO, Inches(0.6), Inches(3.0), width=Inches(3.3))
add_text(s, 5.0, 2.5, 7.8, 1.0, "Plan IA", size=54, bold=True, color=BLANCO)
add_text(s, 5.0, 3.4, 7.8, 0.8, "Operaciones y Finanzas", size=32, color=AZUL)
line = s.shapes.add_connector(1, Inches(5.0), Inches(4.4), Inches(8.5), Inches(4.4))
line.line.color.rgb = AZUL; line.line.width = Pt(2)
add_text(s, 5.0, 4.6, 7.8, 0.4, "Hoja de ruta 2026", size=18, color=BLANCO)
add_text(s, 5.0, 6.0, 7.8, 0.4, "Andrés Browne", size=14, bold=True, color=BLANCO)
add_text(s, 5.0, 6.4, 7.8, 0.4, "Gerencia Finanzas + Supply Chain", size=12, color=GRIS_CLARO)
add_text(s, 5.0, 6.8, 7.8, 0.4, "26 de mayo de 2026", size=12, color=GRIS_CLARO)


# ========== SLIDE 2 — RESUMEN EJECUTIVO ==========
s = add_slide()
add_title(s, "Resumen ejecutivo",
          "El Cerebro IA ya está construido a medias. Cerrando 5 piezas liberamos eficiencias significativas.")
add_rect(s, 0.6, 1.85, 6.0, 3.5, fill=FONDO, line=GRIS_CLARO)
add_rect(s, 6.7, 1.85, 6.0, 3.5, fill=AZUL)
add_text(s, 0.8, 2.0, 5.6, 0.4, "HOY", size=14, bold=True, color=ROJO)
add_text(s, 6.9, 2.0, 5.6, 0.4, "DICIEMBRE 2026", size=14, bold=True, color=BLANCO)
hoy = ["• Procesos manuales · 75% errores", "• Decisiones con dato de hace 30 días",
       "• Excel paralelos en cada PC", "• Equipo captura datos"]
futuro = ["• Procesos automatizados · < 5% errores", "• 3 reportes auto-generados lunes 9:00",
          "• Decisiones con dato fresco (< 1 día)", "• 1 cerebro centralizado consultable",
          "• Equipo valida y analiza"]
add_text(s, 0.8, 2.5, 5.6, 2.7, hoy, size=13, color=GRIS_OSCURO)
add_text(s, 6.9, 2.5, 5.6, 2.7, futuro, size=13, color=BLANCO)
add_text(s, 0.6, 5.7, 12.2, 0.4, "LOS 2 OBJETIVOS DEL PLAN", size=12, bold=True, color=GRIS_MEDIO)
add_rounded(s, 0.6, 6.1, 5.95, 0.95, fill=BLANCO, line=AZUL)
add_text(s, 0.85, 6.2, 5.5, 0.35, "1 · CENTRALIZACIÓN", size=11, bold=True, color=AZUL_OSCURO)
add_text(s, 0.85, 6.55, 5.5, 0.5, "Cerebro Union X — una sola fuente de verdad", size=12, color=GRIS_OSCURO)
add_rounded(s, 6.75, 6.1, 5.95, 0.95, fill=BLANCO, line=VERDE)
add_text(s, 7.0, 6.2, 5.5, 0.35, "2 · REDUCCIÓN DE COSTOS", size=11, bold=True, color=VERDE)
add_text(s, 7.0, 6.55, 5.5, 0.5, "~$33 MM en 2026 · ~$100 MM/año run-rate", size=12, bold=True, color=GRIS_OSCURO)
add_footer(s, 2)


# ========== SLIDE 3 — FOTO 360 ==========
s = add_slide()
add_title(s, "Foto 360 — Lo que ya está vivo HOY",
          "4 apps web · 3 agentes IA · 6 skills · integraciones múltiples")
cards = [
    ("4 APPS WEB", AZUL, ["• App Ventas (17 vistas, 5 usuarios)", "• App Finanzas (P&L, EBIT, FCST)",
                          "• App Operaciones (COMEX, Stock, Bonos)", "• App Planificación (Triada)"]),
    ("3 AGENTES IA", VERDE, ["• Agente COMEX Gmail ✓", "• Agente LCV-Compras ✓", "• Agente Cobranza ⏸ pausado"]),
    ("6 SKILLS CLAUDE", AZUL_OSCURO, ["• comex-workflow · shipping-plan", "• distribución comisiones · reporte gerencial",
                                       "• segmentación pedido · EERR clasificador"]),
    ("INTEGRACIONES", NARANJA, ["• Odoo (XML-RPC)", "• Gmail · Drive · Sheets API",
                                 "• SII (Playwright) · DuckDB (post 7-jun)"]),
]
card_w = 2.95
for i, (title, color, items) in enumerate(cards):
    x = 0.5 + i * (card_w + 0.15)
    add_rounded(s, x, 1.8, card_w, 2.85, fill=BLANCO, line=GRIS_CLARO)
    add_rect(s, x, 1.8, card_w, 0.5, fill=color)
    add_text(s, x + 0.15, 1.88, card_w - 0.3, 0.4, title, size=12, bold=True, color=BLANCO)
    add_text(s, x + 0.15, 2.45, card_w - 0.3, 2.2, items, size=10.5, color=GRIS_OSCURO)

add_text(s, 0.5, 4.85, 12.5, 0.4, "Estado por pieza", size=14, bold=True, color=AZUL_OSCURO)
pieces = [("App Ventas", "Producción", VERDE), ("App Finanzas", "Producción", VERDE),
          ("App Operaciones", "Producción", VERDE), ("App Planificación", "Producción (F2 OK)", VERDE),
          ("Agente COMEX", "Activo polling 2 min", VERDE), ("Agente LCV-Compras", "Daily email", VERDE),
          ("Agente Cobranza", "Pausado · incidente", ROJO), ("Skills Claude", "On-demand", VERDE)]
col_w = 3.0
for i, (name, status, color) in enumerate(pieces):
    col = i % 4; row = i // 4
    x = 0.5 + col * (col_w + 0.05); y = 5.3 + row * 0.75
    add_rounded(s, x, y, col_w, 0.65, fill=BLANCO, line=GRIS_CLARO)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.15), Inches(y + 0.22), Inches(0.18), Inches(0.18))
    dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background(); dot.shadow.inherit = False
    add_text(s, x + 0.45, y + 0.05, col_w - 0.5, 0.3, name, size=11, bold=True, color=GRIS_OSCURO)
    add_text(s, x + 0.45, y + 0.32, col_w - 0.5, 0.3, status, size=9, color=GRIS_MEDIO)
add_footer(s, 3)


# ========== SLIDE 4 — CEREBRO DEFINICIÓN ==========
s = add_slide()
add_title(s, "El Cerebro Union X",
          "Centralizar para tener data en vivo, capturar eficiencias y dejar el valor agregado dentro de la empresa.")
add_text(s, 0.5, 1.8, 6, 0.4, "EL PROBLEMA HOY", size=13, bold=True, color=ROJO)
add_rounded(s, 0.5, 2.2, 6.0, 4.3, fill=FONDO, line=GRIS_CLARO)
problemas = [("Odoo", "ERP potente, pero solo 1 de las fuentes"),
             ("Drives por persona", "Formatos propios · archivos locales"),
             ("Integraciones a medias", "Cada área toma lo que necesita"),
             ("Marketplaces", "ML, Falabella, Paris, Ripley… separados")]
for i, (label, desc) in enumerate(problemas):
    y = 2.4 + i * 1.0
    add_rect(s, 0.75, y, 0.5, 0.5, fill=ROJO)
    add_text(s, 0.75, y, 0.5, 0.5, str(i+1), size=18, bold=True, color=BLANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 1.4, y + 0.0, 5.0, 0.3, label, size=13, bold=True, color=GRIS_OSCURO)
    add_text(s, 1.4, y + 0.3, 5.0, 0.3, desc, size=11, color=GRIS_MEDIO)
add_text(s, 0.75, 6.2, 5.5, 0.3, "→ Cada área arma su versión · valor a medias", size=11, bold=True, color=ROJO)

add_text(s, 6.8, 1.8, 6, 0.4, "LA RESPUESTA — 5 FASES", size=13, bold=True, color=VERDE)
fases = [("F1", "Auditar", "Mapear cómo trabaja cada área"),
         ("F2", "Integrar", "Modelo único semi-automático"),
         ("F3", "Automatizar", "Apagar Drives/Excel paralelos"),
         ("F4", "Eficientar", "Optimizar integraciones"),
         ("F5", "Agentes", "IA replica los flujos del Cerebro")]
for i, (code, title, desc) in enumerate(fases):
    y = 2.2 + i * 0.85
    add_rounded(s, 6.8, y, 5.95, 0.78, fill=BLANCO, line=AZUL)
    add_rect(s, 6.95, y + 0.13, 0.55, 0.5, fill=AZUL)
    add_text(s, 6.95, y + 0.13, 0.55, 0.5, code, size=15, bold=True, color=BLANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 7.7, y + 0.1, 5.0, 0.3, title, size=13, bold=True, color=AZUL_OSCURO)
    add_text(s, 7.7, y + 0.4, 5.0, 0.3, desc, size=11, color=GRIS_MEDIO)
add_footer(s, 4)


# ========== SLIDE 5 — CASO RAW VENTAS ==========
s = add_slide()
add_title(s, "Caso paradigmático: RAW de Ventas",
          "Lo que logramos con Ventas (F1→F3 ya cerrado) lo replicamos en cada área.")
fases5 = [("F1", "Auditar", "Construido por tercero, sin documentar", "41 columnas mapeadas, fuente por fuente", VERDE),
          ("F2", "Integrar", "Info de Odoo + Excel + Drive + MP, manual", "Modelo unificado en Turso, semi-auto", VERDE),
          ("F3", "Automatizar", "Cada área hacía SU versión del RAW", "RAW único en vivo · 17 vistas · 5 usuarios", VERDE),
          ("F4", "Eficientar", "—", "Pendiente: comisiones / logística / mkt por canal", AMARILLO),
          ("F5", "Agentes", "—", "Pendiente: agente auto-cura Matriz + Maestra", ROJO)]
add_rect(s, 0.5, 1.9, 1.0, 0.5, fill=AZUL_OSCURO)
add_text(s, 0.5, 1.95, 1.0, 0.4, "Fase", size=12, bold=True, color=BLANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, 1.55, 1.9, 5.4, 0.5, fill=AZUL_OSCURO)
add_text(s, 1.65, 1.95, 5.3, 0.4, "ANTES", size=12, bold=True, color=BLANCO, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, 7.0, 1.9, 5.4, 0.5, fill=AZUL_OSCURO)
add_text(s, 7.1, 1.95, 5.3, 0.4, "HOY", size=12, bold=True, color=BLANCO, anchor=MSO_ANCHOR.MIDDLE)
for i, (code, title, antes, hoy, color) in enumerate(fases5):
    y = 2.5 + i * 0.85
    add_rect(s, 0.5, y, 1.0, 0.8, fill=color)
    add_text(s, 0.5, y + 0.1, 1.0, 0.3, code, size=14, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    add_text(s, 0.5, y + 0.42, 1.0, 0.3, title, size=10, color=BLANCO, align=PP_ALIGN.CENTER)
    add_rect(s, 1.55, y, 5.4, 0.8, fill=BLANCO, line=GRIS_CLARO)
    add_text(s, 1.7, y + 0.2, 5.2, 0.5, antes, size=11, color=GRIS_OSCURO, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, 7.0, y, 5.4, 0.8, fill=BLANCO, line=GRIS_CLARO)
    add_text(s, 7.15, y + 0.2, 5.2, 0.5, hoy, size=11, color=GRIS_OSCURO, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, 0.5, 6.95, 12.5, 0.3,
         "Mensaje a Gerencias: el patrón está validado. La pregunta no es \"¿se puede?\", sino \"¿cuándo lo replicamos en cada área?\".",
         size=12, bold=True, color=AZUL_OSCURO, align=PP_ALIGN.CENTER)
add_footer(s, 5)


# ========== SLIDE 6 — MIRADA EMPRESA ==========
s = add_slide()
add_title(s, "Mirada Empresa — Avance por área en las 5 fases del Cerebro",
          "9 áreas mapeadas. Cada celda muestra el % de avance.")
areas = [("Ventas",                [100, 100, 100, 30, 0]),
         ("Finanzas + Ctrl Gest.", [100, 100, 70, 10, 0]),
         ("Operaciones",           [100, 100, 60, 0, 30]),
         ("Planificación",         [100, 100, 0, 0, 0]),
         ("Contabilidad (LCV)",    [100, 50, 30, 0, 40]),
         ("SAC / Log. Inv.",       [30, 0, 0, 0, 0]),
         ("Facturación",           [100, 0, 10, 0, 0]),
         ("EEFF",                  [50, 0, 0, 0, 0]),
         ("Productos",             [0, 0, 0, 0, 0])]
add_text(s, 0.5, 1.85, 2.8, 0.4, "Área", size=12, bold=True, color=AZUL_OSCURO)
fase_names = ["F1 Auditar", "F2 Integrar", "F3 Automatizar", "F4 Eficientar", "F5 Agentes"]
col_w = 1.85
for j, name in enumerate(fase_names):
    x = 3.4 + j * col_w
    add_rect(s, x, 1.85, col_w - 0.1, 0.5, fill=AZUL_OSCURO)
    add_text(s, x, 1.9, col_w - 0.1, 0.4, name, size=11, bold=True, color=BLANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
row_h = 0.47
for i, (area, pcts) in enumerate(areas):
    y = 2.5 + i * row_h
    add_text(s, 0.5, y + 0.05, 2.8, row_h - 0.05, area, size=11, bold=True, color=GRIS_OSCURO, anchor=MSO_ANCHOR.MIDDLE)
    for j, pct in enumerate(pcts):
        x = 3.4 + j * col_w
        if pct >= 100: color = VERDE; texto = "✓ 100%"
        elif pct >= 50: color = AMARILLO; texto = f"{pct}%"
        elif pct > 0: color = NARANJA; texto = f"{pct}%"
        else: color = ROJO; texto = "0%"
        add_rect(s, x, y, col_w - 0.1, row_h - 0.05, fill=color)
        add_text(s, x, y, col_w - 0.1, row_h - 0.05, texto, size=11, bold=True, color=BLANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, 0.5, 6.85, 12.5, 0.3,
         "Verde: completo (100%)   ·   Amarillo: avanzado (50-99%)   ·   Naranja: en progreso (1-49%)   ·   Rojo: no iniciado",
         size=10, color=GRIS_MEDIO, align=PP_ALIGN.CENTER)
add_footer(s, 6)


# ========== HELPER ÁREA ==========
def slide_area(num, area, subtitle, fases_data, hito_fecha, hito_desc, hito_color=ROJO):
    s = add_slide()
    add_title(s, area, subtitle)
    col_w = 2.45; gap = 0.1
    for i, (code, title, pct, color, detail) in enumerate(fases_data):
        x = 0.5 + i * (col_w + gap)
        add_rounded(s, x, 1.85, col_w, 4.2, fill=BLANCO, line=GRIS_CLARO)
        add_rect(s, x, 1.85, col_w, 0.7, fill=color)
        add_text(s, x + 0.1, 1.9, col_w - 0.2, 0.3, code, size=14, bold=True, color=BLANCO)
        add_text(s, x + 0.1, 2.2, col_w - 0.2, 0.3, title, size=11, color=BLANCO)
        add_text(s, x + 0.1, 2.7, col_w - 0.2, 0.5, f"{pct}%", size=26, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.1, 3.25, col_w - 0.2, 0.3, "completo" if pct >= 100 else "próximo paso:",
                 size=10, color=GRIS_MEDIO, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.15, 3.6, col_w - 0.3, 2.4, detail, size=10, color=GRIS_OSCURO)
    add_rounded(s, 0.5, 6.2, 12.3, 0.85, fill=hito_color)
    add_text(s, 0.7, 6.3, 1.8, 0.65, "🎯 HITO", size=13, bold=True, color=BLANCO, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, 2.5, 6.3, 0.04, 0.65, fill=BLANCO)
    add_text(s, 2.7, 6.3, 2.2, 0.3, hito_fecha, size=12, bold=True, color=BLANCO)
    add_text(s, 2.7, 6.6, 9.9, 0.4, hito_desc, size=11, color=BLANCO)
    add_footer(s, num)


slide_area(7, "Ventas — Caso paradigmático",
           "El recorrido completo F1→F3 ya está logrado. Ahora F4 (margen real) y F5 (agente curador).",
           [("F1", "Auditar", 100, VERDE, "41 columnas mapeadas, fuente por fuente"),
            ("F2", "Integrar", 100, VERDE, "Modelo unificado en Turso · sync 3×/día"),
            ("F3", "Automatizar", 100, VERDE, "App Ventas única · 17 vistas · 5 usuarios"),
            ("F4", "Eficientar", 30, NARANJA, "Cargar comisiones / logística / mkt por canal. Contenido: Subgcia. Comercial. Arquitectura: Andrés"),
            ("F5", "Agentes", 0, ROJO, "Agente auto-cura Matriz Productos y Maestra Canales · alerta SKUs nuevos")],
           "JUNIO 2026", "Cargar comisiones · logística · marketing por canal → margen final REAL. Subgcia. Comercial toma contenido; arquitectura sigue con Andrés.",
           AZUL_OSCURO)

slide_area(8, "Finanzas + Control de Gestión",
           "App ya integrada. Skills de contribución y P&L se centralizan en junio.",
           [("F1", "Auditar", 100, VERDE, "EERR clasificado · 88 reglas · mapeo cuentas analíticas. Control Gestión: scope definido"),
            ("F2", "Integrar", 100, VERDE, "App Finanzas: P&L 7 líneas · EBIT consolidado · FCST Ene-Dic"),
            ("F3", "Automatizar", 70, AMARILLO, "Sigue dependiendo del P&L sheet + cálculos P&L y contribución por Drive"),
            ("F4", "Eficientar", 10, ROJO, "Incorporar archivos MD de contribución y P&L al cerebro"),
            ("F5", "Agentes", 0, ROJO, "Todo centralizado · EEFF desde Odoo en vivo · agente construye reportes")],
           "JUNIO 2026", "Skills de Contribución y P&L quedan en el cerebro. Pendiente: EEFF integrado.",
           AZUL_OSCURO)

slide_area(9, "Operaciones — COMEX · Stock · Fulfillment",
           "Fulfillment 100% auto excepto costos operativos. COMEX tránsito ok · falta pre-embarque y producción.",
           [("F1", "Auditar", 100, VERDE, "Flujos COMEX (PI→PL→OHNSO) · stock LIVE · capacidad WMS"),
            ("F2", "Integrar", 100, VERDE, "Apps integradas · Stock LIVE desde Odoo · Fulfillment 100% auto (todo Odoo)"),
            ("F3", "Automatizar", 60, NARANJA, "Falta automatizar costos operativos · COMEX pre-embarque y producción"),
            ("F4", "Eficientar", 0, ROJO, "Eficiencia: Slotting · rotación inventario · prever costos variables · punto equilibrio vs planificación venta · oportunidades anticipadas"),
            ("F5", "Agentes", 30, NARANJA, "COMEX activo polling 2 min. Próximo: COMEX v2 sin validación intermedia")],
           "JUNIO 2026", "Automatización COMEX pre-embarque + producción → conecta con Triada de Planificación.",
           AZUL_OSCURO)

slide_area(10, "Contabilidad — Libro de Compra y Venta",
           "Agente LCV-Compras en producción. Cobranza pausado por incidente — fix prioridad #1.",
           [("F1", "Auditar", 100, VERDE, "Mapeo SII × Odoo × Drive completo (compra y venta)"),
            ("F2", "Integrar", 50, AMARILLO, "LCV-Compras: 3 fuentes cruzadas diario. Próximo: LCV-Venta + Cobranza"),
            ("F3", "Automatizar", 30, NARANJA, "LCV-Compras: cuadre 3-vías + email 7 AM. Próximo: fix incidente Cobranza"),
            ("F4", "Eficientar", 0, ROJO, "Validación cruzada con auditoría externa automatizada"),
            ("F5", "Agentes", 40, NARANJA, "Agente LCV-Compras en prod. Cobranza pausado. Próximo: fix + 2 ciclos validados")],
           "31 JUL · OCT 2026", "Salidas de Analista Contable #1 (31/07) y Analista Contable #2 (Oct) habilitadas por LCV + Cobranza estables.",
           ROJO)

slide_area(11, "SAC / Logística Inversa",
           "Solo F1 al 30%. Acelerar F1→F3 antes del 31-jul para liberar Analista Log. Inversa.",
           [("F1", "Auditar", 30, NARANJA, "URGENTE: completar auditoría flujos NC + devoluciones"),
            ("F2", "Integrar", 0, ROJO, "Modelo unificado NC en Turso · integración con WMS"),
            ("F3", "Automatizar", 0, ROJO, "Informe semanal 100% auto · email sin intervención"),
            ("F4", "Eficientar", 0, ROJO, "Pendiente F1-F3"),
            ("F5", "Agentes", 0, ROJO, "Pendiente F1-F3")],
           "31 JULIO 2026", "Salida de Analista Logística Inversa. El informe semanal queda 100% automatizado.",
           ROJO)

slide_area(12, "Facturación",
           "POC etiqueta + deadline activo. Cerrar F2-F3 antes de 15-ago y 15-sep.",
           [("F1", "Auditar", 100, VERDE, "Flujos por canal mapeados · criterios facturación documentados"),
            ("F2", "Integrar", 0, ROJO, "Centralizar criterio facturación por canal en tabla única"),
            ("F3", "Automatizar", 10, ROJO, "POC etiqueta + deadline activo. Próximo: producción + 2 ciclos validados"),
            ("F4", "Eficientar", 0, ROJO, "Pendiente F2-F3"),
            ("F5", "Agentes", 0, ROJO, "Pendiente F2-F3")],
           "15 AGO · 15 SEP 2026", "Salida de Facturadora #1 (15/08) y Facturadora #2 (15/09). POC etiqueta+deadline las reemplaza.",
           ROJO)

slide_area(13, "EEFF — Estados Financieros",
           "EEFF en vivo desde Odoo con reglas contables, integrados a la planilla de planificación financiera.",
           [("F1", "Auditar", 50, AMARILLO, "Próximo: completar mapeo cuentas faltantes · reglas contables documentadas"),
            ("F2", "Integrar", 0, ROJO, "Conectar Odoo → modelo EEFF en vivo aplicando reglas contables"),
            ("F3", "Automatizar", 0, ROJO, "EEFF en vivo en planilla planif. financiera · idealmente automatizada en la app"),
            ("F4", "Eficientar", 0, ROJO, "Pendiente F2-F3"),
            ("F5", "Agentes", 0, ROJO, "Pendiente F2-F3")],
           "30 SEPTIEMBRE 2026", "Salida de Control de Gestión (Gabriela). EEFF en vivo en planilla planif. financiera.",
           ROJO)

slide_area(14, "Planificación — Triada (stock + llegadas + demanda)",
           "App ya en producción F2. En junio el equipo de Planificación toma 100% el control del contenido.",
           [("F1", "Auditar", 100, VERDE, "Flujos planificación auditados · Triada definida: stock + llegadas + demanda → cobertura"),
            ("F2", "Integrar", 100, VERDE, "App Planificación en producción · Triada Proyectada (baseline + live)"),
            ("F3", "Automatizar", 0, ROJO, "Equipo Planificación toma control del contenido · feed COMEX pre-embarque y producción"),
            ("F4", "Eficientar", 0, ROJO, "Pendiente F3"),
            ("F5", "Agentes", 0, ROJO, "Agente sugiere ajustes a planificación según señales de demanda real")],
           "JUNIO 2026", "Equipo Planificación toma 100% el control del contenido. La Triada se alimenta de COMEX pre-embarque + producción.",
           AZUL_OSCURO)


# ========== SLIDE 15 — EFICIENCIAS OVERVIEW (directos) ==========
s = add_slide()
add_title(s, "Eficiencias directas identificadas",
          "Ahorros calculados con cifras del P&L. Eficiencia en GAV operacional para 2027.")
add_rounded(s, 0.5, 1.85, 6.1, 2.0, fill=AZUL)
add_text(s, 0.7, 2.0, 5.8, 0.4, "AHORRO 2026", size=14, bold=True, color=BLANCO)
add_text(s, 0.7, 2.4, 5.8, 1.0, "$32,85 MM", size=58, bold=True, color=BLANCO)
add_text(s, 0.7, 3.4, 5.8, 0.4, "junio-diciembre 2026", size=12, color=BLANCO)
add_rounded(s, 6.7, 1.85, 6.1, 2.0, fill=VERDE)
add_text(s, 6.9, 2.0, 5.8, 0.4, "EFICIENCIA GAV OPERACIONAL 2027", size=12, bold=True, color=BLANCO)
add_text(s, 6.9, 2.4, 5.8, 1.0, "$99,7 MM", size=58, bold=True, color=BLANCO)
add_text(s, 6.9, 3.4, 5.8, 0.4, "run-rate anualizado · solo costos directos", size=12, color=BLANCO)
add_text(s, 0.5, 4.1, 12.5, 0.4, "DESGLOSE POR CATEGORÍA", size=13, bold=True, color=GRIS_MEDIO)
cats = [("Personal (6 cargos + cargas)", "$28,4 MM", "$91,0 MM", AZUL_OSCURO, "91%"),
        ("Reducción usuarios Odoo (35 → 15)", "$2,7 MM", "$5,47 MM", AZUL, "5%"),
        ("Eliminación Multivende (desde jun)", "$1,15 MM", "$1,98 MM", AZUL, "2%"),
        ("Fulfillment Falabella +10%", "$0,6 MM", "$1,28 MM", AZUL, "1%")]
add_rect(s, 0.5, 4.6, 6.5, 0.5, fill=AZUL_OSCURO)
add_text(s, 0.7, 4.65, 6.3, 0.4, "Categoría", size=11, bold=True, color=BLANCO, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, 7.05, 4.6, 2.0, 0.5, fill=AZUL_OSCURO)
add_text(s, 7.05, 4.65, 2.0, 0.4, "2026", size=11, bold=True, color=BLANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, 9.1, 4.6, 2.0, 0.5, fill=AZUL_OSCURO)
add_text(s, 9.1, 4.65, 2.0, 0.4, "Run-rate", size=11, bold=True, color=BLANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, 11.15, 4.6, 1.65, 0.5, fill=AZUL_OSCURO)
add_text(s, 11.15, 4.65, 1.65, 0.4, "% del total", size=11, bold=True, color=BLANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
for i, (cat, m26, mrun, color, pct) in enumerate(cats):
    y = 5.15 + i * 0.5
    bg = FONDO if i % 2 == 0 else BLANCO
    add_rect(s, 0.5, y, 6.5, 0.45, fill=bg, line=GRIS_CLARO)
    add_text(s, 0.7, y, 6.3, 0.45, cat, size=11, color=GRIS_OSCURO, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, 7.05, y, 2.0, 0.45, fill=bg, line=GRIS_CLARO)
    add_text(s, 7.05, y, 2.0, 0.45, m26, size=11, bold=True, color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, 9.1, y, 2.0, 0.45, fill=bg, line=GRIS_CLARO)
    add_text(s, 9.1, y, 2.0, 0.45, mrun, size=11, bold=True, color=VERDE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, 11.15, y, 1.65, 0.45, fill=bg, line=GRIS_CLARO)
    add_text(s, 11.15, y, 1.65, 0.45, pct, size=11, color=GRIS_MEDIO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 15)


# ========== SLIDE 16 — REDUCCIÓN PLANTILLA ==========
s = add_slide()
add_title(s, "Reducción de plantilla por automatización",
          "6 cargos liberados entre julio y octubre 2026. Transición planificada con IA en producción.")
cargos = [("Analista Logística Inversa", "31/07/2026", "$1.004.852", "5 m", "$5,02 MM", "$12,06 MM"),
          ("Analista Contable #1", "31/07/2026", "$1.015.754", "5 m", "$5,08 MM", "$12,19 MM"),
          ("Facturadora #1", "15/08/2026", "$895.815", "4,5 m", "$4,03 MM", "$10,75 MM"),
          ("Facturadora #2", "15/09/2026", "$822.212", "3,5 m", "$2,88 MM", "$9,87 MM"),
          ("Control de Gestión", "30/09/2026", "$2.258.096", "3 m", "$6,77 MM", "$27,10 MM"),
          ("Analista Contable #2", "Oct 2026", "$770.000", "2 m", "$1,54 MM", "$9,24 MM")]
headers = ["Cargo", "Sale", "Bruto/mes", "Meses 2026", "Ahorro 2026", "Anualizado"]
cols_x = [0.5, 4.5, 6.0, 7.55, 8.8, 10.55]
cols_w = [3.95, 1.45, 1.5, 1.2, 1.7, 2.25]
for i, h in enumerate(headers):
    add_rect(s, cols_x[i], 1.85, cols_w[i] - 0.05, 0.5, fill=AZUL_OSCURO)
    add_text(s, cols_x[i] + 0.1, 1.9, cols_w[i] - 0.25, 0.4, h, size=11, bold=True, color=BLANCO,
             align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
for i, row in enumerate(cargos):
    y = 2.4 + i * 0.5
    bg = FONDO if i % 2 == 0 else BLANCO
    for j, val in enumerate(row):
        add_rect(s, cols_x[j], y, cols_w[j] - 0.05, 0.45, fill=bg, line=GRIS_CLARO)
        color = AZUL_OSCURO if j == 4 else (VERDE if j == 5 else GRIS_OSCURO)
        bold = (j >= 4)
        add_text(s, cols_x[j] + 0.1, y, cols_w[j] - 0.25, 0.45, val, size=10.5, bold=bold, color=color,
                 align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
y_tot = 2.4 + 6 * 0.5
add_rect(s, 0.5, y_tot, 7.55, 0.5, fill=AZUL_OSCURO)
add_text(s, 0.7, y_tot, 7.3, 0.5, "TOTAL (con cargas patronales 12%)", size=12, bold=True, color=BLANCO, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, 8.05, y_tot, 2.45, 0.5, fill=AZUL_OSCURO)
add_text(s, 8.05, y_tot, 2.45, 0.5, "~$28,4 MM", size=13, bold=True, color=BLANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, 10.55, y_tot, 2.25, 0.5, fill=VERDE)
add_text(s, 10.55, y_tot, 2.25, 0.5, "~$91,0 MM/año", size=13, bold=True, color=BLANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, 0.5, 6.0, 12.5, 0.4, "LÍNEA DE TIEMPO DE SALIDAS", size=12, bold=True, color=GRIS_MEDIO)
line = s.shapes.add_connector(1, Inches(1.0), Inches(6.85), Inches(12.5), Inches(6.85))
line.line.color.rgb = AZUL; line.line.width = Pt(2)
months = ["JUL", "AGO", "SEP", "OCT"]
for i, m in enumerate(months):
    x = 1.5 + i * 2.85
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - 0.2), Inches(6.65), Inches(0.4), Inches(0.4))
    circle.fill.solid(); circle.fill.fore_color.rgb = AZUL; circle.line.fill.background(); circle.shadow.inherit = False
    add_text(s, x - 0.5, 6.65, 1.0, 0.4, m, size=10, bold=True, color=BLANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, 1.0, 6.3, 1.5, 0.3, "31 Jul", size=9, bold=True, color=ROJO)
add_text(s, 1.0, 7.1, 1.8, 0.3, "Anal.LI + AC#1", size=9, color=GRIS_OSCURO)
add_text(s, 3.85, 6.3, 1.5, 0.3, "15 Ago", size=9, bold=True, color=NARANJA)
add_text(s, 3.85, 7.1, 1.5, 0.3, "Facturadora 1", size=9, color=GRIS_OSCURO)
add_text(s, 6.7, 6.3, 2.5, 0.3, "15 Sep · 30 Sep", size=9, bold=True, color=NARANJA)
add_text(s, 6.7, 7.1, 2.5, 0.3, "Facturadora 2 + Ctrl Gestión", size=9, color=GRIS_OSCURO)
add_text(s, 9.55, 6.3, 1.5, 0.3, "Oct", size=9, bold=True, color=NARANJA)
add_text(s, 9.55, 7.1, 1.8, 0.3, "Anal. Contable #2", size=9, color=GRIS_OSCURO)
add_footer(s, 16)


# ========== SLIDE 17 — PASO A PASO REDUCCIÓN ==========
s = add_slide()
add_title(s, "Paso a paso para la reducción",
          "Cada salida requiere IA en producción y validada 2 meses antes. Transición planificada.")
headers = ["Cargo", "IA que reemplaza", "Estado IA", "Prerrequisito", "Fecha"]
cols_x = [0.5, 3.05, 5.6, 7.05, 11.65]
cols_w = [2.5, 2.5, 1.4, 4.55, 1.15]
for i, h in enumerate(headers):
    add_rect(s, cols_x[i], 1.85, cols_w[i] - 0.05, 0.5, fill=AZUL_OSCURO)
    add_text(s, cols_x[i] + 0.08, 1.9, cols_w[i] - 0.2, 0.4, h, size=11, bold=True, color=BLANCO, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
rows = [("Analista Log. Inv.", "Agente SAC / Log. Inversa", "F1 30%", ROJO,
         "Auditar flujos NC · automatizar informe semanal · 2 ciclos validados", "31/07"),
        ("Analista Contable #1", "Agente LCV-Compras (activo)", "F5 50%", VERDE,
         "Onboarding Camila/Víctor · agente estable 2 meses", "31/07"),
        ("Analista Contable #2", "Agente Cobranza", "F3 0%", ROJO,
         "Fix root cause incidente · 2 ciclos validados · Camila absorbe carga", "Oct"),
        ("Facturadora #1", "POC etiqueta + deadline", "F3 10%", AMARILLO,
         "POC a producción · 2 ciclos sin error · validación Yohana", "15/08"),
        ("Facturadora #2", "POC + flujo #1 validado", "Sigue #1", AMARILLO,
         "Flujo #1 estable post-15-ago · Yohana redistribuye carga", "15/09"),
        ("Control de Gestión", "Cerebro + EEFF en vivo + Chat ejec.", "F4 10%", AMARILLO,
         "EEFF en vivo desde Odoo · planilla planif. financiera automatizada · chat ejec. POC", "30/09")]
for i, (cargo, ia, estado, color_est, prereq, fecha) in enumerate(rows):
    y = 2.4 + i * 0.7
    bg = FONDO if i % 2 == 0 else BLANCO
    add_rect(s, cols_x[0], y, cols_w[0] - 0.05, 0.65, fill=bg, line=GRIS_CLARO)
    add_text(s, cols_x[0] + 0.1, y, cols_w[0] - 0.25, 0.65, cargo, size=10.5, bold=True, color=GRIS_OSCURO, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, cols_x[1], y, cols_w[1] - 0.05, 0.65, fill=bg, line=GRIS_CLARO)
    add_text(s, cols_x[1] + 0.1, y, cols_w[1] - 0.25, 0.65, ia, size=10, color=GRIS_OSCURO, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, cols_x[2], y, cols_w[2] - 0.05, 0.65, fill=bg, line=GRIS_CLARO)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cols_x[2] + 0.1), Inches(y + 0.22), Inches(0.18), Inches(0.18))
    dot.fill.solid(); dot.fill.fore_color.rgb = color_est; dot.line.fill.background(); dot.shadow.inherit = False
    add_text(s, cols_x[2] + 0.35, y, cols_w[2] - 0.4, 0.65, estado, size=10, bold=True, color=color_est, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, cols_x[3], y, cols_w[3] - 0.05, 0.65, fill=bg, line=GRIS_CLARO)
    add_text(s, cols_x[3] + 0.1, y, cols_w[3] - 0.25, 0.65, prereq, size=9.5, color=GRIS_OSCURO, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, cols_x[4], y, cols_w[4] - 0.05, 0.65, fill=bg, line=GRIS_CLARO)
    add_text(s, cols_x[4] + 0.1, y, cols_w[4] - 0.25, 0.65, fecha, size=11, bold=True, color=AZUL_OSCURO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rounded(s, 0.5, 6.65, 12.3, 0.45, fill=AZUL)
add_text(s, 0.7, 6.7, 11.9, 0.35,
         "💡 Cada salida requiere IA en producción 2 meses antes con validación humana cruzada · Transición planificada.",
         size=11, bold=True, color=BLANCO, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 17)


# ========== SLIDE 18 — OTRAS EFICIENCIAS DIRECTAS ==========
s = add_slide()
add_title(s, "Otras eficiencias directas habilitadas por IA",
          "3 palancas adicionales que la centralización del Cerebro destraba.")
cards = [("Reducción usuarios Odoo", "35 → 15", "$5,47 MM/año",
          "Apps centralizadas (Ventas + Finanzas + Ops + Planificación) reemplazan acceso directo de 20 personas.",
          "Audit usos 90d · capacitación · apps cubren flujos", AZUL_OSCURO),
         ("Eliminación Multivende", "DESDE JUNIO", "$1,98 MM/año",
          "Integraciones directas Odoo ↔ marketplaces vía Agente COMEX reemplazan Multivende.",
          "Definir integración alternativa · 1 ciclo validado · rollback plan", AZUL),
         ("Fulfillment Falabella +10%", "+10% share", "$1,28 MM/año",
          "Tarifa fulfillment Fala ($2.490/ped) < costo bodega allocated ($2.887/ped) · +$397/ped a favor.",
          "Costo variable real · plan migración SKUs · decisión Comité", VERDE)]
card_w = 4.05
for i, (title, sub, ahorro, desc, prereq, color) in enumerate(cards):
    x = 0.5 + i * (card_w + 0.15)
    add_rounded(s, x, 1.85, card_w, 5.2, fill=BLANCO, line=GRIS_CLARO)
    add_rect(s, x, 1.85, card_w, 0.75, fill=color)
    add_text(s, x + 0.2, 1.95, card_w - 0.4, 0.55, title, size=14, bold=True, color=BLANCO, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + 0.2, 2.85, card_w - 0.4, 0.5, sub, size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.2, 3.45, card_w - 0.4, 0.5, ahorro, size=22, bold=True, color=VERDE, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.25, 4.1, card_w - 0.5, 1.5, desc, size=11, color=GRIS_OSCURO)
    add_rect(s, x + 0.25, 5.55, card_w - 0.5, 0.05, fill=GRIS_CLARO)
    add_text(s, x + 0.25, 5.65, card_w - 0.5, 0.3, "PRERREQUISITO", size=9, bold=True, color=GRIS_MEDIO)
    add_text(s, x + 0.25, 5.95, card_w - 0.5, 1.0, prereq, size=10.5, color=GRIS_OSCURO)
add_footer(s, 18)


# ========== SLIDE 19 — AHORROS POTENCIALES (NUEVO) ==========
s = add_slide()
add_title(s, "Ahorros POTENCIALES — escenario conservador",
          "Más allá de los costos directos: lo que la IA habilita en eficiencia operativa y comercial.")

# Cifra grande arriba
add_rounded(s, 0.5, 1.85, 12.3, 1.6, fill=VERDE_OSCURO)
add_text(s, 0.7, 2.0, 11.9, 0.4, "POTENCIAL ADICIONAL · RUN-RATE ANUAL", size=13, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
add_text(s, 0.7, 2.4, 11.9, 0.9, "+$58 MM/año", size=58, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
add_text(s, 0.7, 3.15, 11.9, 0.3, "ESCENARIO CONSERVADOR — depende de implementación correcta", size=11, color=BLANCO, align=PP_ALIGN.CENTER)

# 3 palancas
palancas = [
    ("📊 FORECASTING\nReducir quiebres", "$28 MM/año",
     "Tasa quiebres actual estimada: ~5% demanda perdida (multi-canal)",
     "Mejor forecasting → reducir a 3% (-2 pp)",
     "Recupero 2% × $5.000 MM venta = $100 MM en ventas",
     "Margen 56% = $56 MM · Conservador 50% = $28 MM"),
    ("📦 INVENTARIO\n-0,5 mes", "$15 MM/año",
     "Inventario promedio estimado: ~$250 MM",
     "Reducir 0,5 mes = $125 MM capital trabajo liberado",
     "Costo financiero (12% anual) = $15 MM/año",
     "+ Costo operativo evitado ~$3 MM/año"),
    ("🎯 MARGEN FINAL\nMix óptimo canales", "$15 MM/año",
     "Hoy margen_final = margen_front (sin comisión/log/mkt)",
     "Con margen real → cortar canales tóxicos",
     "Mejora margen +1 pp sobre $5.000 MM = $50 MM",
     "Conservador 30% = $15 MM/año"),
]
card_w = 4.05
for i, (titulo, cifra, l1, l2, l3, l4) in enumerate(palancas):
    x = 0.5 + i * (card_w + 0.15)
    add_rounded(s, x, 3.65, card_w, 3.35, fill=BLANCO, line=GRIS_CLARO)
    add_rect(s, x, 3.65, card_w, 0.85, fill=VERDE)
    add_text(s, x + 0.15, 3.7, card_w - 0.3, 0.8, titulo, size=12, bold=True, color=BLANCO, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + 0.15, 4.65, card_w - 0.3, 0.5, cifra, size=24, bold=True, color=VERDE, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.25, 5.25, card_w - 0.5, 1.7, [l1, l2, l3, l4], size=10, color=GRIS_OSCURO)
add_footer(s, 19)


# ========== SLIDE 20 — QUÉ NECESITAMOS ==========
s = add_slide()
add_title(s, "Qué necesitamos para lograr las otras eficiencias",
          "Prerrequisitos técnicos + organizacionales con owner y deadline.")
headers = ["Eficiencia", "Prerrequisito técnico", "Prerrequisito organizacional", "Owner", "Deadline"]
cols_x = [0.5, 3.0, 6.5, 10.05, 11.75]
cols_w = [2.45, 3.45, 3.5, 1.65, 1.1]
for i, h in enumerate(headers):
    add_rect(s, cols_x[i], 1.85, cols_w[i] - 0.05, 0.5, fill=AZUL_OSCURO)
    add_text(s, cols_x[i] + 0.1, 1.9, cols_w[i] - 0.2, 0.4, h, size=11, bold=True, color=BLANCO, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
rows = [("Eliminación Multivende",
         "Integración directa Odoo ↔ MP via Agente COMEX · 1 ciclo validado",
         "Validación Felipe · plan rollback si falla",
         "Andrés + Felipe", "JUN 2026"),
        ("Fulfillment +10% Fala",
         "Costo variable evitable real · plan migración SKUs · stock fulfillment",
         "Decisión Comité Comercial · política inventario",
         "Andrés + Comité", "Q4 2026"),
        ("Odoo 35 → 15 usuarios",
         "Apps cubren flujos de 20 usuarios · capacitación 2 sesiones · audit 90d",
         "Comunicación al equipo · acuerdos por área",
         "Andrés + IT", "Q3 2026")]
for i, row in enumerate(rows):
    y = 2.4 + i * 1.4
    bg = FONDO if i % 2 == 0 else BLANCO
    for j, val in enumerate(row):
        add_rect(s, cols_x[j], y, cols_w[j] - 0.05, 1.35, fill=bg, line=GRIS_CLARO)
        color = AZUL_OSCURO if j == 0 else GRIS_OSCURO
        bold = (j == 0) or (j == 4)
        align = PP_ALIGN.CENTER if j == 4 else PP_ALIGN.LEFT
        add_text(s, cols_x[j] + 0.1, y + 0.05, cols_w[j] - 0.25, 1.25, val, size=10.5, bold=bold, color=color, align=align, anchor=MSO_ANCHOR.MIDDLE)
add_rounded(s, 0.5, 6.65, 12.3, 0.45, fill=NARANJA)
add_text(s, 0.7, 6.7, 11.9, 0.35,
         "⚠️ Esta tabla es el insumo principal de la Carta Gantt del siguiente bloque.",
         size=11, bold=True, color=BLANCO, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 20)


# ========== SLIDE 21 — GANTT MACRO ==========
s = add_slide()
add_title(s, "Carta Gantt — Vista macro",
          "3 capas: Cerebro (por área) · Salidas personal · Otras eficiencias.")
meses = ["Jun", "Jul", "Ago", "Sep", "Oct", "Nov", "Dic"]
start_x = 4.0; month_w = 1.2
for i, m in enumerate(meses):
    x = start_x + i * month_w
    add_rect(s, x, 1.85, month_w - 0.02, 0.4, fill=AZUL_OSCURO)
    add_text(s, x, 1.85, month_w - 0.02, 0.4, m, size=11, bold=True, color=BLANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, 0.5, 2.35, 3.4, 0.4, "Hitos personal", size=10, bold=True, color=ROJO)
real_hito_xs = [(5.2 + 1.2 * 1, "31 Jul"), (6.4 + 1.2 * 0.5, "15 Ago"),
                (7.6 + 1.2 * 0.5, "15 Sep"), (7.6 + 1.2 * 1, "30 Sep"),
                (8.8 + 1.2 * 0.5, "Oct")]
for x_pos, label in real_hito_xs:
    tri = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x_pos - 0.15), Inches(2.35), Inches(0.3), Inches(0.35))
    tri.fill.solid(); tri.fill.fore_color.rgb = ROJO; tri.line.fill.background(); tri.shadow.inherit = False
tracks = [("SAC F1→F3", 0, 1.5, NARANJA),
          ("Contabilidad LCV F2→F3", 0, 1.5, NARANJA),
          ("Facturación F2→F3", 0.5, 2.0, AMARILLO),
          ("EEFF F1→F3", 1.5, 3.5, AMARILLO),
          ("Finanzas F3→F4", 0, 3.0, AZUL),
          ("Ventas F4 Mg Real", 0, 2.0, AZUL),
          ("Operaciones COMEX v2 + pre-embarque", 0, 4.0, AZUL),
          ("Planificación F3 (control equipo)", 0, 1.0, AZUL),
          ("Cerebro Voz (POC+prod)", 3.5, 6.5, AZUL_OSCURO)]
efs = [("Multivende eliminado", 0, 1.0, VERDE),
       ("Odoo 35→15", 1.0, 3.0, VERDE),
       ("Fulfillment Fala +10%", 3.5, 6.5, VERDE)]
add_text(s, 0.5, 2.85, 3.4, 0.3, "Cerebro por área", size=10, bold=True, color=AZUL_OSCURO)
for i, (name, mstart, mend, color) in enumerate(tracks):
    y = 3.15 + i * 0.3
    add_text(s, 0.5, y, 3.4, 0.3, name, size=9, color=GRIS_OSCURO, anchor=MSO_ANCHOR.MIDDLE)
    bar_x = start_x + mstart * month_w
    bar_w = (mend - mstart) * month_w
    add_rect(s, bar_x, y + 0.04, bar_w - 0.05, 0.2, fill=color)
line2 = s.shapes.add_connector(1, Inches(0.5), Inches(6.0), Inches(12.5), Inches(6.0))
line2.line.color.rgb = GRIS_CLARO; line2.line.width = Pt(0.75)
add_text(s, 0.5, 6.05, 3.4, 0.3, "Otras eficiencias", size=10, bold=True, color=VERDE)
for i, (name, mstart, mend, color) in enumerate(efs):
    y = 6.4 + i * 0.3
    add_text(s, 0.5, y, 3.4, 0.3, name, size=10, color=GRIS_OSCURO, anchor=MSO_ANCHOR.MIDDLE)
    bar_x = start_x + mstart * month_w
    bar_w = (mend - mstart) * month_w
    add_rect(s, bar_x, y + 0.03, bar_w - 0.05, 0.22, fill=color)
add_footer(s, 21)


# ========== SLIDE 22 — HITOS CRÍTICOS ==========
s = add_slide()
add_title(s, "Hitos críticos — Entregables previos a cada salida",
          "Para cada fecha, los entregables IA tienen que estar en producción y validados antes.")
hitos = [("31 JUL", ROJO, "Anal. Log. Inv. + Anal. Cont. #1",
          ["Agente SAC F1 completo (15-jun)", "Agente SAC F2-F3 informe auto (30-jun)",
           "2 ciclos validados sin error (15-jul)", "Agente LCV-Compras 60d prod (15-jul)",
           "Onboarding Camila/Víctor (20-jul)"]),
         ("15 AGO", NARANJA, "Facturadora #1",
          ["POC etiqueta+deadline en prod (30-jun)", "2 ciclos completos sin error (julio)",
           "Tabla criterio facturación canal (15-jul)", "Validación final flujo (10-ago)"]),
         ("15 SEP", NARANJA, "Facturadora #2",
          ["Flujo #1 estable post-15-ago", "Redistribución carga residual",
           "Validación 30 días con flujo #1 sola"]),
         ("30 SEP", ROJO, "Control de Gestión",
          ["EEFF F1 completo (31-jul)", "Skill reporte integrada al pipeline (15-ago)",
           "EEFF en vivo + planilla planif. financiera (15-sep)", "Skills contribución y P&L en cerebro",
           "Chat ejecutivo POC operativo (22-sep)"]),
         ("OCT", AMARILLO, "Anal. Contable #2",
          ["Fix root cause Agente Cobranza (15-jun urgente)",
           "Cobranza reactivada + 60 d sin error (15-sep)", "Camila absorbe carga residual"])]
col_w = 2.45
for i, (fecha, color, salida, items) in enumerate(hitos):
    x = 0.5 + i * (col_w + 0.05)
    add_rect(s, x, 1.85, col_w, 0.7, fill=color)
    add_text(s, x + 0.1, 1.9, col_w - 0.2, 0.3, fecha, size=14, bold=True, color=BLANCO)
    add_text(s, x + 0.1, 2.2, col_w - 0.2, 0.3, salida, size=9, color=BLANCO)
    add_rounded(s, x, 2.55, col_w, 4.45, fill=BLANCO, line=GRIS_CLARO)
    text_lines = [f"• {it}" for it in items]
    add_text(s, x + 0.15, 2.7, col_w - 0.3, 4.2, text_lines, size=9.5, color=GRIS_OSCURO)
add_footer(s, 22)


# ========== SLIDE 23 — CALENDARIO ==========
s = add_slide()
add_title(s, "Vista calendario — Qué se entrega cada mes",
          "Cinco meses críticos jun→oct 2026.")
months_data = [("JUNIO", AZUL, ["SAC F1-F3", "Cobranza fix root cause", "Comisión canal Ventas",
                                "Skills Contribución+P&L cerebro", "Planificación → equipo",
                                "COMEX pre-embarque + producción", "Multivende eliminado"], None),
               ("JULIO", AZUL_OSCURO, ["SAC F3 validación", "POC facturación", "LCV-Compras 60d",
                                       "Onboarding Camila/Víctor"], ("31 Jul", "Anal.LI + AC #1")),
               ("AGOSTO", NARANJA, ["Fact #1 en prod", "Odoo audit usos 90d",
                                    "Costo variable bodega"], ("15 Ago", "Fact. #1")),
               ("SEPTIEMBRE", NARANJA, ["Fact #2 validación", "EEFF en vivo + planilla",
                                        "Skills P&L 2 ciclos", "Chat ejecutivo POC"], ("15+30 Sep", "Fact #2 + Ctrl Gest.")),
               ("OCTUBRE", AMARILLO, ["Cobranza reactivada", "Cerebro Voz POC",
                                      "Fulfillment plan migración"], ("Oct", "Anal. Cont. #2"))]
col_w = 2.45
for i, (mes, color, items, hito) in enumerate(months_data):
    x = 0.5 + i * (col_w + 0.05)
    add_rect(s, x, 1.85, col_w, 0.55, fill=color)
    add_text(s, x, 1.85, col_w, 0.55, mes, size=14, bold=True, color=BLANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rounded(s, x, 2.4, col_w, 3.5, fill=BLANCO, line=GRIS_CLARO)
    items_lines = [f"• {it}" for it in items]
    add_text(s, x + 0.15, 2.55, col_w - 0.3, 3.3, items_lines, size=10, color=GRIS_OSCURO)
    if hito:
        add_rect(s, x, 5.95, col_w, 1.05, fill=ROJO if color in [ROJO, AMARILLO] else NARANJA)
        add_text(s, x + 0.1, 6.0, col_w - 0.2, 0.4, f"▼ {hito[0]}", size=12, bold=True, color=BLANCO)
        add_text(s, x + 0.1, 6.4, col_w - 0.2, 0.55, hito[1], size=10, color=BLANCO, anchor=MSO_ANCHOR.MIDDLE)
    else:
        add_rect(s, x, 5.95, col_w, 1.05, fill=FONDO, line=GRIS_CLARO)
        add_text(s, x, 5.95, col_w, 1.05, "—", size=14, color=GRIS_MEDIO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 23)


# ========== SLIDE 24 — KPIs + BENCHMARK ==========
s = add_slide()
add_title(s, "KPIs — Cuadro de mando dic-2026",
          "4 dimensiones · Benchmark sector: ROI digital típico 2-3x año 1 · 4-6x año 2 exitoso.")
cards = [("COSTOS", AZUL_OSCURO, [("Headcount admin", "47 → 41", "(-6 personas)"),
                                   ("Costo personal/año", "-$91 MM", "run-rate"),
                                   ("Costo SaaS", "-$8,7 MM", "Odoo+Multivende"),
                                   ("Inversión IA", "USD $2.250/m", "Claude + Cerebro")]),
         ("EFICIENCIA", AZUL, [("Reportes manuales/sem", "5 → 0", "Andrés"),
                                ("Tiempo consolidación", "↓ 80%", "≤ 2 hrs/sem"),
                                ("Lead time COMEX", "↓ 70%", "≤ 1 día"),
                                ("Visibilidad EBIT", "30 d → 1 d", "")]),
         ("CALIDAD", VERDE, [("Tasa error reportes", "75% → < 5%", ""),
                              ("Cuadre SII×Odoo", "5-7 d → < 1 d", ""),
                              ("Quiebres demanda", "5% → 3%", "potencial"),
                              ("Inventario", "-0,5 mes", "potencial")]),
         ("MADUREZ + ROI", NARANJA, [("Áreas F3+ (de 9)", "3 → 7", ""),
                                       ("Agentes en prod", "2 → 5", ""),
                                       ("ROI 2026", "~3,0x", "directos"),
                                       ("ROI 2027 + potenciales", "~6,6x", "$158 MM/año")])]
card_w = 3.0
for i, (title, color, items) in enumerate(cards):
    x = 0.5 + i * (card_w + 0.1)
    add_rounded(s, x, 1.85, card_w, 5.0, fill=BLANCO, line=GRIS_CLARO)
    add_rect(s, x, 1.85, card_w, 0.55, fill=color)
    add_text(s, x, 1.85, card_w, 0.55, title, size=13, bold=True, color=BLANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for j, (kpi, val, note) in enumerate(items):
        if not kpi: continue
        y = 2.6 + j * 1.05
        add_text(s, x + 0.15, y, card_w - 0.3, 0.3, kpi, size=10, color=GRIS_MEDIO)
        add_text(s, x + 0.15, y + 0.3, card_w - 0.3, 0.45, val, size=16, bold=True, color=color)
        if note:
            add_text(s, x + 0.15, y + 0.75, card_w - 0.3, 0.25, note, size=9, color=GRIS_MEDIO)
add_footer(s, 24)


# ========== SLIDE 25 — RIESGOS ==========
s = add_slide()
add_title(s, "Riesgos y mitigaciones",
          "8 riesgos identificados. El #1 (fix Agente Cobranza) es bloqueante y prioridad máxima.")
headers = ["#", "Riesgo", "P", "I", "Mitigación", "Owner"]
cols_x = [0.5, 1.0, 5.5, 6.05, 6.6, 11.6]
cols_w = [0.5, 4.5, 0.55, 0.55, 5.0, 1.25]
for i, h in enumerate(headers):
    add_rect(s, cols_x[i], 1.85, cols_w[i] - 0.05, 0.4, fill=AZUL_OSCURO)
    add_text(s, cols_x[i] + 0.05, 1.85, cols_w[i] - 0.1, 0.4, h, size=10, bold=True, color=BLANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
riesgos = [("1", "Agente Cobranza no se arregla en junio → caen 2 hitos (jul + oct)", "M", "A",
            "Fix prioridad máxima S1-S2 (1-30 jun) · Plan B: Martín reactiva Task Scheduler local",
            "Andrés + Víctor", ROJO),
           ("2", "Adopción baja del equipo (siguen con Excel local)", "M", "A",
            "Onboarding formal · cerrar acceso planillas paralelas · sponsor Gerencia",
            "Andrés + Gerencia", NARANJA),
           ("3", "Comunicación interna mal manejada → conflicto laboral", "M", "A",
            "Plan comunicación RR.HH. 30 d antes · narrativa automatización · finiquitos OK",
            "Gerencia + RR.HH.", NARANJA),
           ("4", "POC etiqueta+deadline no llega a prod a tiempo (15-ago)", "M", "M",
            "Sprint dedicado junio · Yohana valida semana a semana · plan B: postergar 30-ago",
            "Andrés + Yohana", AMARILLO),
           ("5", "Costo Claude API se descontrola", "B", "M",
            "Prompt caching · límites por skill · auditoría mensual · alerta a USD $3.000",
            "Andrés", AMARILLO),
           ("6", "Incidente tipo Cobranza se repite", "B", "A",
            "Regla nueva: features que escriben en Drive/Odoo requieren OK humano pre-merge",
            "Andrés", AMARILLO),
           ("7", "Andrés es single point of failure", "A", "A",
            "Docs vivos (docs/*_ESTADO_*.md) · onboarding técnico a 2do · Claude copiloto",
            "Andrés", AMARILLO),
           ("8", "Dependencia Odoo (si cae, todo cae)", "B", "A",
            "Histórico parquet · apps funcionan con histórico aunque Odoo caiga",
            "—", VERDE)]
for i, (num, riesgo, p, im, mitig, owner, color) in enumerate(riesgos):
    y = 2.3 + i * 0.58
    bg = FONDO if i % 2 == 0 else BLANCO
    add_rect(s, cols_x[0], y, cols_w[0] - 0.05, 0.55, fill=color)
    add_text(s, cols_x[0] + 0.05, y, cols_w[0] - 0.1, 0.55, num, size=12, bold=True, color=BLANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, cols_x[1], y, cols_w[1] - 0.05, 0.55, fill=bg, line=GRIS_CLARO)
    add_text(s, cols_x[1] + 0.1, y, cols_w[1] - 0.2, 0.55, riesgo, size=9.5, color=GRIS_OSCURO, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, cols_x[2], y, cols_w[2] - 0.05, 0.55, fill=bg, line=GRIS_CLARO)
    add_text(s, cols_x[2], y, cols_w[2] - 0.05, 0.55, p, size=11, bold=True, color=GRIS_OSCURO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, cols_x[3], y, cols_w[3] - 0.05, 0.55, fill=bg, line=GRIS_CLARO)
    add_text(s, cols_x[3], y, cols_w[3] - 0.05, 0.55, im, size=11, bold=True, color=GRIS_OSCURO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, cols_x[4], y, cols_w[4] - 0.05, 0.55, fill=bg, line=GRIS_CLARO)
    add_text(s, cols_x[4] + 0.1, y, cols_w[4] - 0.2, 0.55, mitig, size=9.5, color=GRIS_OSCURO, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, cols_x[5], y, cols_w[5] - 0.05, 0.55, fill=bg, line=GRIS_CLARO)
    add_text(s, cols_x[5] + 0.05, y, cols_w[5] - 0.1, 0.55, owner, size=9, color=GRIS_OSCURO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, 0.5, 6.95, 12.5, 0.3, "P = Probabilidad · I = Impacto · A = Alto · M = Medio · B = Bajo",
         size=9, color=GRIS_MEDIO, align=PP_ALIGN.CENTER)
add_footer(s, 25)


# ========== SLIDE 26 — CIERRE ==========
s = add_slide()
add_title(s, "Cierre — Decisiones que necesitamos hoy",
          "5 decisiones que destraban la ejecución de los próximos 30 días.")
decisiones = [("1", "Aprobación del plan global", "Cerebro + 6 salidas + 3 eficiencias directas + potenciales", "Gerencia", "Hoy"),
              ("2", "Aprobación presupuesto IA", "Subir USD $1.500/m → ~$2.250/m (chat ejec. + 5 agentes nuevos)", "Finanzas", "Esta semana"),
              ("3", "Validación calendario salidas", "RR.HH. + Legal", "Gerencia + RR.HH.", "Esta semana"),
              ("4", "Sponsor para adopción", "Cerrar Excel paralelos", "Gerencia", "Inmediato"),
              ("5", "Fix Agente Cobranza prioridad #1", "Destrabar recursos", "Andrés", "30-jun")]
for i, (num, accion, detalle, quien, cuando) in enumerate(decisiones):
    y = 1.85 + i * 0.65
    add_rounded(s, 0.5, y, 12.3, 0.55, fill=BLANCO, line=GRIS_CLARO)
    add_rect(s, 0.5, y, 0.6, 0.55, fill=AZUL_OSCURO)
    add_text(s, 0.5, y, 0.6, 0.55, num, size=18, bold=True, color=BLANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 1.2, y + 0.08, 4.8, 0.2, accion, size=12, bold=True, color=GRIS_OSCURO)
    add_text(s, 1.2, y + 0.3, 4.8, 0.25, detalle, size=10, color=GRIS_MEDIO)
    add_text(s, 6.5, y, 3.5, 0.55, quien, size=11, color=AZUL_OSCURO, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 10.5, y, 2.0, 0.55, cuando, size=11, bold=True, color=ROJO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, 0.5, 5.4, 12.5, 0.4, "PRÓXIMOS 30 DÍAS (JUNIO 2026)", size=12, bold=True, color=GRIS_MEDIO)
sprints = [("S1 (1-7)", "Kick-off + comunicación interna"),
           ("S2 (8-14)", "Fix Agente Cobranza + SAC F1"),
           ("S3 (15-21)", "POC etiqueta + comisiones canal + COMEX pre-embarque"),
           ("S4 (22-30)", "Validación + go/no-go salidas julio")]
for i, (sem, desc) in enumerate(sprints):
    x = 0.5 + i * 3.1
    add_rounded(s, x, 5.85, 3.0, 0.6, fill=AZUL)
    add_text(s, x + 0.15, 5.9, 2.8, 0.25, sem, size=10, bold=True, color=BLANCO)
    add_text(s, x + 0.15, 6.15, 2.8, 0.3, desc, size=9.5, color=BLANCO)
add_rounded(s, 0.5, 6.6, 12.3, 0.5, fill=VERDE)
add_text(s, 0.7, 6.65, 11.9, 0.4,
         "💡 El Cerebro UnionX ya existe a medias. Cerrar las piezas que faltan libera ~$33 MM este año, ~$100 MM run-rate directos + ~$58 MM potenciales.",
         size=11, bold=True, color=BLANCO, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 26)


# ========== SLIDE 27 — ANEXO TÉCNICO ==========
s = add_slide()
add_title(s, "Anexo técnico — Stack · Costo IA · ROI completo",
          "Desglose de inversión + 3 escenarios de ROI (2026 · 2027 directos · 2027 con potenciales).")

# IZQUIERDA — Desglose costo IA
add_text(s, 0.5, 1.85, 7.0, 0.4, "DESGLOSE COSTO IA MENSUAL", size=13, bold=True, color=AZUL_OSCURO)
stack_detail = [("Claude API hoy", "$1.500", "2 agentes + 6 skills + desarrollo"),
                ("+ Agente SAC", "+$80", "Polling NC/devoluciones"),
                ("+ Agente Cobranza (post-fix)", "+$80", "Procesamiento diario 5 clientes"),
                ("+ Agente Facturación", "+$100", "POC etiqueta+deadline prod"),
                ("+ Agente EEFF", "+$100", "EEFF en vivo desde Odoo"),
                ("+ Chat ejecutivo 24/7", "+$400", "Pregúntale al Cerebro"),
                ("+ Alertas tiempo real", "+$150", "10 alertas activas"),
                ("− Prompt caching", "−$200", "Cachear contexto repetido")]
for i, (item, cost, desc) in enumerate(stack_detail):
    y = 2.3 + i * 0.45
    add_rounded(s, 0.5, y, 7.0, 0.4, fill=BLANCO, line=GRIS_CLARO)
    add_text(s, 0.65, y, 3.8, 0.4, item, size=10, bold=True, color=GRIS_OSCURO, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 4.55, y, 1.2, 0.4, cost, size=11, bold=True, color=AZUL_OSCURO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 5.85, y, 1.55, 0.4, desc, size=9, color=GRIS_MEDIO, anchor=MSO_ANCHOR.MIDDLE)

add_rounded(s, 0.5, 5.95, 7.0, 0.5, fill=AZUL)
add_text(s, 0.65, 5.95, 4.0, 0.5, "TOTAL RÉGIMEN COMPLETO", size=12, bold=True, color=BLANCO, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, 4.55, 5.95, 1.2, 0.5, "~$2.250/m", size=14, bold=True, color=BLANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, 5.85, 5.95, 1.55, 0.5, "$27 MM CLP/año", size=10, color=BLANCO, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, 0.5, 6.55, 7.0, 0.3, "+ DuckDB sobre parquet local: $0/mes (post 7-jun)", size=10, color=GRIS_MEDIO)

# DERECHA — ROI 3 escenarios
add_text(s, 7.8, 1.85, 5.2, 0.4, "ROI EN 3 ESCENARIOS", size=13, bold=True, color=AZUL_OSCURO)

roi_data = [("2026", "directos", "$32,85 MM", "$11 MM", "~3,0x", NARANJA),
            ("2027", "directos run-rate", "$99,7 MM", "$24 MM", "~4,2x", VERDE),
            ("2027", "+ potenciales", "$157,7 MM", "$24 MM", "~6,6x", VERDE_OSCURO)]
for i, (year, tipo, ahorro, costo, roi, color) in enumerate(roi_data):
    y = 2.3 + i * 1.5
    add_rounded(s, 7.8, y, 5.2, 1.35, fill=BLANCO, line=color)
    add_rect(s, 7.8, y, 5.2, 0.4, fill=color)
    add_text(s, 7.95, y, 5.0, 0.4, f"{year} · {tipo}", size=11, bold=True, color=BLANCO, anchor=MSO_ANCHOR.MIDDLE)
    # ROI grande
    add_text(s, 7.95, y + 0.45, 1.8, 0.85, roi, size=36, bold=True, color=color, align=PP_ALIGN.CENTER)
    # Detalles
    add_text(s, 9.8, y + 0.5, 3.1, 0.3, "Ahorro:", size=9, color=GRIS_MEDIO)
    add_text(s, 9.8, y + 0.7, 3.1, 0.3, ahorro, size=12, bold=True, color=GRIS_OSCURO)
    add_text(s, 9.8, y + 1.0, 3.1, 0.3, f"Costo IA: {costo}", size=10, color=GRIS_MEDIO)

# Nota benchmark
add_rounded(s, 7.8, 6.85, 5.2, 0.3, fill=FONDO, line=GRIS_CLARO)
add_text(s, 7.95, 6.85, 5.0, 0.3, "Benchmark sector: 2-3x año 1 · 4-6x año 2 exitoso",
         size=9, bold=True, color=AZUL_OSCURO, anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 27)


# ========== GUARDAR ==========
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"OK: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
